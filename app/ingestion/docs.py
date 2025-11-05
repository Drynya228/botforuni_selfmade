# app/ingestion/docs.py
from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from app.ingestion.ocr import ocr_photo

async def extract_text_from_path(path: str) -> str:
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            reader = PdfReader(path)
            parts = [(page.extract_text() or "") for page in reader.pages]
            text = "\n".join(parts).strip()
            return text or await ocr_photo(path)
        if ext in {".jpg", ".jpeg", ".png"}:
            return await ocr_photo(path)
        if ext == ".docx":
            doc = Document(path)
            return "\n".join(par.text for par in doc.paragraphs).strip()
        if ext == ".xlsx":
            wb = load_workbook(path, read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        out.append(" \t".join(vals))
            return "\n".join(out).strip()
        if ext == ".pptx":
            prs = Presentation(path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        out.append(shape.text)
            return "\n".join(out).strip()
    except Exception:
        pass
    return ""
