from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from app.ingestion.ocr import ocr_photo

async def extract_document_text(path: str) -> str:
    p = Path(path)
    ext = p.suffix.lower()
    if ext == ".pdf":
        try:
            txt = []
            reader = PdfReader(path)
            for page in reader.pages:
                txt.append(page.extract_text() or "")
            text = "\n".join(txt).strip()
            if text:
                return text
        except Exception:  # noqa: BLE001
            pass
        # fallback: OCR первой страницы скрина
        return await ocr_photo(path)
    if ext in {".jpg", ".jpeg", ".png"}:
        return await ocr_photo(path)
    if ext == ".docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs).strip()
    if ext == ".xlsx":
        wb = load_workbook(path, read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    out.append(" \t".join(vals))
        return "\n".join(out).strip()
    if ext == ".pptx":
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text)
        return "\n".join(out).strip()
    return ""